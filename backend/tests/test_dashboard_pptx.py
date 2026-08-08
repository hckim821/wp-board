"""대시보드 PPTX 내보내기 — plan.md §0.5.6.

세 가지를 본다.

1. **정말 열리는 파일인가.** 200 과 content-type 만 보면 깨진 바이트도 통과한다.
   그래서 응답을 python-pptx 로 **되읽어** 도형과 텍스트를 단언한다.
2. **레이아웃 규칙** — 순수 함수라 DB 없이 돈다. 미배정이 맨 뒤인지, 조각난
   Phase 가 컬럼 하나로 합쳐지는지, 한 장에 들어가는지.
3. **색 상수가 프론트와 어긋나지 않는지.** §0.5 가 정본이고 값이 두 곳에 적혀
   있으므로, `frontend/src/theme/dashboard.ts` 를 직접 파싱해 대조한다.
   (프론트 파일은 **읽기만** 한다.)
"""

from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path
from urllib.parse import quote

import pytest

from app.models import Item, ItemOwner, ItemStatus
from app.schemas.item import ItemOut, OwnerRef
from app.services import dashboard_theme
from app.schemas.item import DocumentRef
from app.services.dashboard_pptx_service import (
    ROWS_PER_SLIDE,
    build_dashboard_pptx,
    build_detail_pages,
    build_layout,
    compute_metrics,
)

API = "/api/v1"

REPO_ROOT = Path(__file__).resolve().parents[2]
FRONTEND_THEME = REPO_ROOT / "frontend" / "src" / "theme" / "dashboard.ts"


# =============================================================================
# 순수 레이아웃 — DB 없이
# =============================================================================
def row(
    no: int,
    *,
    phase_id: int | None = None,
    phase_no: int | None = None,
    phase_name: str = "",
    milestone_id: int | None = None,
    milestone_no: int | None = None,
    milestone_name: str = "",
    status: ItemStatus = ItemStatus.NOT_STARTED,
    dash_label: str | None = None,
    deliverable: str | None = None,
    title: str | None = None,
    owners: list[str] | None = None,
    documents: list[str] | None = None,
) -> ItemOut:
    return ItemOut(
        id=no,
        sort_order=no,
        row_no=no,
        phase_id=phase_id,
        phase_no=phase_no,
        phase_name=phase_name,
        milestone_id=milestone_id,
        milestone_no=milestone_no,
        milestone_name=milestone_name,
        milestone_no_display=(
            f"{phase_no}.{milestone_no}"
            if phase_no is not None and milestone_no is not None
            else None
        ),
        documents=[
            DocumentRef(id=i, no=int(c), name=f"문서 {c}")
            for i, c in enumerate(documents or [], start=1)
        ],
        is_phase_block_start=False,
        is_phase_block_end=False,
        is_milestone_block_start=False,
        is_milestone_block_end=False,
        can_create_phase=False,
        can_create_milestone=False,
        status=status,
        dash_label=dash_label,
        deliverable=deliverable,
        title=title,
        owners=[OwnerRef(id=i, name=n) for i, n in enumerate(owners or [], start=1)],
    )


def test_layout_groups_phases_then_milestones():
    bands = build_layout([
        row(1, phase_id=1, phase_no=0, phase_name="Pre", milestone_id=10, milestone_no=1,
            milestone_name="Gap"),
        row(2, phase_id=1, phase_no=0, phase_name="Pre", milestone_id=10, milestone_no=1,
            milestone_name="Gap"),
        row(3, phase_id=1, phase_no=0, phase_name="Pre", milestone_id=11, milestone_no=2,
            milestone_name="I/O"),
        row(4, phase_id=2, phase_no=1, phase_name="Init", milestone_id=20, milestone_no=1,
            milestone_name="Scope"),
    ])

    assert [b.label for b in bands] == ["Phase 0. Pre", "Phase 1. Init"]
    assert [len(b.columns) for b in bands] == [2, 1]
    assert [len(c.cards) for c in bands[0].columns] == [2, 1]
    assert bands[0].columns[0].number_label == "0.1"
    assert bands[0].columns[1].name == "I/O"


def test_unassigned_rows_become_a_trailing_colourless_band():
    bands = build_layout([
        row(1, phase_id=1, phase_no=0, phase_name="Pre", milestone_id=10, milestone_no=1),
        row(2),                                   # 회색 행
        row(3, phase_id=2, phase_no=1, phase_name="Init", milestone_id=20, milestone_no=1),
    ])

    assert [b.label for b in bands] == ["Phase 0. Pre", "Phase 1. Init", "미배정"]
    assert bands[-1].color == dashboard_theme.DASH_UNASSIGNED
    assert [c.no for c in bands[-1].columns[0].cards] == [2]


def test_no_unassigned_band_when_every_row_is_assigned():
    bands = build_layout([
        row(1, phase_id=1, phase_no=0, phase_name="Pre", milestone_id=10, milestone_no=1)
    ])
    assert [b.label for b in bands] == ["Phase 0. Pre"]


def test_a_fragmented_phase_stays_one_column():
    """조각난 보드를 같은 이름표의 컬럼 둘로 그리지 않는다 — 그건 V4/V5 가 할 말이다."""
    bands = build_layout([
        row(1, phase_id=1, phase_no=0, phase_name="Pre", milestone_id=10, milestone_no=1),
        row(2, phase_id=2, phase_no=1, phase_name="Init", milestone_id=20, milestone_no=1),
        row(3, phase_id=1, phase_no=0, phase_name="Pre", milestone_id=10, milestone_no=1),
    ])

    assert [b.label for b in bands] == ["Phase 0. Pre", "Phase 1. Init"]
    assert [c.no for c in bands[0].columns[0].cards] == [1, 3]


def test_a_phase_without_a_milestone_gets_a_nameless_bucket():
    """두 계층 모두에서 합법인 상태다 — 떨어뜨리면 행이 사라진다."""
    bands = build_layout([
        row(1, phase_id=1, phase_no=0, phase_name="Pre"),
    ])
    assert bands[0].columns[0].name == "미지정"
    assert bands[0].columns[0].number_label == ""


def test_phase_colours_cycle_past_the_fourth():
    bands = build_layout([
        row(i + 1, phase_id=i + 1, phase_no=i, phase_name=f"P{i}",
            milestone_id=(i + 1) * 10, milestone_no=1)
        for i in range(5)
    ])
    assert bands[4].color == bands[0].color        # 팔레트가 4개라 다섯 번째가 순환


def test_card_text_follows_the_fallback_chain():
    bands = build_layout([
        row(1, phase_id=1, phase_no=0, milestone_id=10, milestone_no=1,
            dash_label="라벨", deliverable="산출물", title="제목"),
        row(2, phase_id=1, phase_no=0, milestone_id=10, milestone_no=1,
            deliverable="산출물만", title="제목"),
        row(3, phase_id=1, phase_no=0, milestone_id=10, milestone_no=1, title="제목만"),
        row(4, phase_id=1, phase_no=0, milestone_id=10, milestone_no=1),
    ])
    assert [c.text for c in bands[0].columns[0].cards] == ["라벨", "산출물만", "제목만", ""]


# =============================================================================
# 한 장에 맞춘다
# =============================================================================
def reference_board() -> list[ItemOut]:
    """기준 케이스 — 엑셀 원본과 같은 Phase 4 / Milestone 13 / 35행.

    Phase 별 행 수 4·10·11·10, Milestone 수 2·3·4·4 (plan.md §1.1).
    """
    shape = {0: [2, 2], 1: [3, 4, 3], 2: [3, 3, 3, 2], 3: [3, 3, 2, 2]}
    items: list[ItemOut] = []
    no = 0
    for phase_seq, milestones in shape.items():
        for ms_index, count in enumerate(milestones, start=1):
            for _ in range(count):
                no += 1
                items.append(
                    row(
                        no,
                        phase_id=phase_seq + 1,
                        phase_no=phase_seq,
                        phase_name=f"Phase{phase_seq}",
                        milestone_id=(phase_seq + 1) * 100 + ms_index,
                        milestone_no=ms_index,
                        milestone_name=f"마일스톤 {phase_seq}.{ms_index}",
                        dash_label=f"라벨 {no}",
                        owners=["DSEP 인프라 담당자"],
                    )
                )
    return items


def test_the_reference_board_is_thirteen_columns():
    bands = build_layout(reference_board())
    assert sum(len(b.columns) for b in bands) == 13
    assert sum(len(c.cards) for b in bands for c in b.columns) == 35


def test_the_status_map_fits_on_one_slide():
    """가로·세로 모두 슬라이드 안에 들어와야 한다."""
    from app.services.dashboard_pptx_service import (
        CARDS_BOTTOM, CARDS_Y, MARGIN_RIGHT, MARGIN_X, SLIDE_W,
    )

    bands = build_layout(reference_board())
    m = compute_metrics(bands)

    total_columns = sum(len(b.columns) for b in bands)
    assert MARGIN_X + total_columns * m.column_pitch <= SLIDE_W - MARGIN_RIGHT + 1e-9

    tallest = max(len(c.cards) for b in bands for c in b.columns)
    assert CARDS_Y + tallest * m.card_pitch <= CARDS_BOTTOM + 1e-9


def test_the_reference_case_keeps_the_source_decks_dimensions():
    """기준 케이스에서는 축소가 걸리지 않아야 한다.

    정본 덱의 피치 0.95 · 카드 0.82 는 **상한**이다. 13 컬럼짜리 실제 보드가
    그 상한에 그대로 앉지 못하면 자동 축소 식이 너무 공격적이라는 뜻이다.
    """
    from app.services.dashboard_pptx_service import CARD_H, CARD_PT, COL_PITCH

    m = compute_metrics(build_layout(reference_board()))
    assert m.column_pitch == pytest.approx(COL_PITCH)
    assert m.card_h == pytest.approx(CARD_H)
    assert m.card_pt == pytest.approx(CARD_PT)


def test_a_denser_board_shrinks_rather_than_overflowing():
    """컬럼이 배로 늘어도 폭이 줄어들 뿐 넘치지 않는다."""
    from app.services.dashboard_pptx_service import MARGIN_RIGHT, MARGIN_X, SLIDE_W

    dense = [
        row(i + 1, phase_id=1, phase_no=0, phase_name="P",
            milestone_id=100 + i, milestone_no=i + 1)
        for i in range(30)
    ]
    m = compute_metrics(build_layout(dense))
    assert MARGIN_X + 30 * m.column_pitch <= SLIDE_W - MARGIN_RIGHT + 1e-9
    assert m.card_pt >= 5.5                       # 폰트에 하한이 있다


def test_a_tall_column_shrinks_the_cards():
    """세로도 마찬가지 — 한 마일스톤에 카드가 몰리면 카드가 낮아진다."""
    from app.services.dashboard_pptx_service import CARDS_BOTTOM, CARDS_Y, CARD_H

    tall = [
        row(i + 1, phase_id=1, phase_no=0, milestone_id=10, milestone_no=1)
        for i in range(20)
    ]
    m = compute_metrics(build_layout(tall))
    assert m.card_h < CARD_H
    assert CARDS_Y + 20 * m.card_pitch <= CARDS_BOTTOM + 1e-9


def test_the_card_height_is_one_value_for_the_whole_board():
    """§0.5.4b — 컬럼마다 다르면 같은 줄의 카드가 어긋나 보인다.

    `Metrics.card_h` 가 스칼라 하나라는 것이 그 성질의 구현이다. 컬럼별 높이를
    도입하려면 이 타입부터 바뀌므로, 이 단언이 그 변경을 막는다.
    """
    lopsided = [
        row(1, phase_id=1, phase_no=0, milestone_id=10, milestone_no=1),
        *[row(i + 2, phase_id=1, phase_no=0, milestone_id=11, milestone_no=2) for i in range(9)],
    ]
    m = compute_metrics(build_layout(lopsided))
    assert isinstance(m.card_h, float) and m.card_h > 0


def test_metrics_survive_an_empty_board():
    m = compute_metrics(build_layout([]))
    assert m.column_w > 0 and m.card_h > 0


# =============================================================================
# 상세 슬라이드 분할 (§0.5.6 개정)
# =============================================================================
def test_each_phase_becomes_its_own_detail_pages():
    pages = build_detail_pages(reference_board())

    # Phase 4개 → 4·10·11·10 행 → 1 + 2 + 3 + 2 = 8장
    assert [p.group_label for p in pages][:1] == ["Phase 0. Phase0"]
    assert len(pages) == 1 + 2 + 3 + 2
    assert all(len(p.rows) <= ROWS_PER_SLIDE for p in pages)


def test_a_phase_of_five_rows_or_fewer_is_a_single_page():
    """경계 바로 아래 — 나뉘지 않으므로 페이지 표기가 붙지 않는다."""
    items = [
        row(i + 1, phase_id=1, phase_no=0, phase_name="P", milestone_id=10, milestone_no=1)
        for i in range(ROWS_PER_SLIDE)
    ]
    pages = build_detail_pages(items)

    assert len(pages) == 1
    assert pages[0].page_count == 1
    assert "(1/" not in pages[0].title
    assert pages[0].title.endswith(f"— 항목 1~{ROWS_PER_SLIDE}")


def test_one_row_past_the_boundary_splits_into_two_pages():
    """경계 바로 위 — 나뉘고, 제목에 `(1/2)` 와 항목 범위가 붙는다."""
    items = [
        row(i + 1, phase_id=1, phase_no=0, phase_name="P", milestone_id=10, milestone_no=1)
        for i in range(ROWS_PER_SLIDE + 1)
    ]
    pages = build_detail_pages(items)

    assert len(pages) == 2
    assert [len(p.rows) for p in pages] == [ROWS_PER_SLIDE, 1]
    assert f"(1/2) — 항목 1~{ROWS_PER_SLIDE}" in pages[0].title
    assert f"(2/2) — 항목 {ROWS_PER_SLIDE + 1}~{ROWS_PER_SLIDE + 1}" in pages[1].title


def test_the_item_range_uses_the_board_row_numbers():
    """`항목 N~M` 은 페이지 내 순번이 아니라 보드의 행 번호다."""
    items = [
        row(i, phase_id=1, phase_no=0, phase_name="P", milestone_id=10, milestone_no=1)
        for i in range(11, 11 + ROWS_PER_SLIDE)
    ]
    assert build_detail_pages(items)[0].title.endswith(
        f"— 항목 11~{10 + ROWS_PER_SLIDE}"
    )


def test_unassigned_rows_get_their_own_group_at_the_end():
    pages = build_detail_pages([
        row(1, phase_id=1, phase_no=0, phase_name="P", milestone_id=10, milestone_no=1),
        row(2),
        row(3),
    ])

    assert [p.group_label for p in pages] == ["Phase 0. P", "미배정"]
    assert pages[-1].eyebrow == "UNASSIGNED ITEMS"
    assert [r.no for r in pages[-1].rows] == [2, 3]


def test_detail_rows_carry_deliverable_milestone_documents_and_owners():
    pages = build_detail_pages([
        row(1, phase_id=1, phase_no=0, phase_name="P", milestone_id=10, milestone_no=1,
            milestone_name="Gap 점검", deliverable="DSEP Gap & Resource Plan",
            title="환경을 점검하고 계획을 확정", documents=["2"],
            owners=["DSEP 인프라 담당자"]),
    ])

    row_out = pages[0].rows[0]
    assert row_out.deliverable == "DSEP Gap & Resource Plan"
    assert row_out.milestone == "0.1 Gap 점검"
    assert row_out.documents == ["2"]
    assert row_out.owners == ["DSEP 인프라 담당자"]


def test_an_empty_board_produces_no_detail_pages():
    assert build_detail_pages([]) == []


# =============================================================================
# 실제 파일 — python-pptx 로 되읽는다
# =============================================================================
def open_pptx(payload: bytes):
    from pptx import Presentation

    return Presentation(BytesIO(payload))


def texts_of(presentation, index: int = 0) -> list[str]:
    return [
        shape.text_frame.text
        for shape in presentation.slides[index].shapes
        if shape.has_text_frame and shape.text_frame.text
    ]


def all_texts(presentation) -> list[str]:
    return [t for i in range(len(presentation.slides)) for t in texts_of(presentation, i)]


def test_the_generated_file_opens_at_16_by_9():
    from pptx.util import Inches

    deck = open_pptx(build_dashboard_pptx("데모 프로젝트", reference_board()))

    assert deck.slide_width == Inches(13.333)
    assert deck.slide_height == Inches(7.5)


def test_the_deck_is_one_status_map_plus_the_detail_pages():
    """슬라이드 수 = 1 + 상세. 실제 파일에서 센다."""
    items = reference_board()
    deck = open_pptx(build_dashboard_pptx("데모 프로젝트", items))
    assert len(deck.slides) == 1 + len(build_detail_pages(items))
    assert len(deck.slides) == 9              # Status Map + 8장


def test_the_status_map_carries_the_title_milestones_and_cards():
    deck = open_pptx(build_dashboard_pptx("데모 프로젝트", reference_board()))
    texts = texts_of(deck, 0)

    assert "PROJECT BOARD  ·  STATUS MAP" in texts
    assert any("데모 프로젝트" in t and "35항목" in t for t in texts)
    assert any(t.startswith("Phase 0.") for t in texts)
    # 마일스톤 헤더는 정본 덱 형식 그대로 — 번호 줄 + 이름 줄.
    assert "0.1\n마일스톤 0.1" in texts
    # 카드는 **두 문단**이다: No 만 있는 줄, 그다음 라벨 줄 (§0.5-2 개정).
    assert "1\n라벨 1" in texts and "35\n라벨 35" in texts


def test_a_card_never_joins_its_number_and_label_on_one_line():
    """`1 | Gap·자원 계획` 처럼 한 줄로 붙으면 번호가 문장의 일부로 읽힌다.

    §0.5-2 개정: index 는 가운데 정렬로 자기 줄을 갖고, 라벨은 그 아래 왼쪽
    정렬이다 — 웹 대시보드 카드와 같은 포맷.
    """
    deck = open_pptx(build_dashboard_pptx("카드", reference_board()))
    slide = deck.slides[0]

    cards = [
        shape.text_frame
        for shape in slide.shapes
        if shape.has_text_frame and re.fullmatch(r"\d+\n라벨 \d+", shape.text_frame.text)
    ]
    assert len(cards) == 35, "카드 35장이 전부 두 문단이어야 한다"

    for frame in cards:
        assert len(frame.paragraphs) == 2
        assert " | " not in frame.text
        head, body = frame.paragraphs
        assert head.text.isdigit()
        assert head.runs[0].font.bold is True
        assert body.runs[0].font.bold is False


def test_card_text_is_anchored_to_the_top_of_the_card():
    """가운데 앵커면 라벨이 한 줄인 카드와 두 줄인 카드의 번호 높이가 달라진다.

    카드 높이는 이미 전 보드에서 하나로 통일돼 있으므로(§0.5.4b), 위를 기준으로
    맞춰야 격자를 훑는 눈이 번호 줄을 잃지 않는다.
    """
    from pptx.enum.text import MSO_ANCHOR

    deck = open_pptx(build_dashboard_pptx("앵커", reference_board()))
    cards = [
        shape.text_frame
        for shape in deck.slides[0].shapes
        if shape.has_text_frame and re.fullmatch(r"\d+\n라벨 \d+", shape.text_frame.text)
    ]
    assert len(cards) == 35
    assert all(f.vertical_anchor == MSO_ANCHOR.TOP for f in cards)


def test_a_card_with_no_label_is_a_single_centred_number():
    """라벨이 비면 둘째 문단을 만들지 않는다 — 빈 줄이 카드를 흔들지 않게."""
    items = [row(1, phase_id=1, phase_no=0, milestone_id=10, milestone_no=1)]
    deck = open_pptx(build_dashboard_pptx("라벨 없음", items))

    frames = [
        shape.text_frame
        for shape in deck.slides[0].shapes
        if shape.has_text_frame and shape.text_frame.text == "1"
    ]
    assert frames, "번호만 있는 카드가 있어야 한다"
    assert all(len(f.paragraphs) == 1 for f in frames)


def test_the_detail_slides_carry_their_rows():
    deck = open_pptx(build_dashboard_pptx("데모 프로젝트", reference_board()))
    second = texts_of(deck, 1)

    assert "PROJECT BOARD  ·  PHASE 0 ITEMS" in second
    assert any("— 항목 1~4" in t for t in second)
    assert "산출물 (Deliverable)  ·  Milestone  /  Key Action" in second
    assert "문서" in second and "Owner" in second
    # 푸터 — 프로젝트명 + 페이지 번호
    assert "데모 프로젝트" in second and "2" in second


def test_every_card_becomes_shapes_on_the_status_map():
    """35장이면 도형이 그만큼 늘어야 한다 — 조용히 빠진 카드가 없는지."""
    few = open_pptx(build_dashboard_pptx("적음", reference_board()[:5]))
    many = open_pptx(build_dashboard_pptx("많음", reference_board()))
    assert len(many.slides[0].shapes) > len(few.slides[0].shapes)


def test_a_split_phase_shows_the_page_marker_on_the_slide():
    """11행짜리 Phase 2 는 3장으로 나뉘고 제목이 그것을 말한다."""
    deck = open_pptx(build_dashboard_pptx("분할", reference_board()))
    titles = [t for t in all_texts(deck) if "(1/" in t or "(2/" in t or "(3/" in t]

    assert any("(1/3)" in t for t in titles)
    assert any("(3/3)" in t for t in titles)


def test_the_detail_rows_show_deliverable_milestone_documents_and_owner():
    items = [
        row(1, phase_id=1, phase_no=0, phase_name="Pre", milestone_id=10, milestone_no=1,
            milestone_name="Gap 점검", deliverable="DSEP Gap & Resource Plan",
            title="환경 점검", documents=["2"], owners=["DSEP 인프라 담당자"]),
    ]
    detail = texts_of(open_pptx(build_dashboard_pptx("한 행", items)), 1)

    assert any("DSEP Gap & Resource Plan" in t and "0.1 Gap 점검" in t for t in detail)
    assert any("환경 점검" in t for t in detail)
    assert "2" in detail
    assert "DSEP 인프라 담당자" in detail
    assert "1" in detail                        # No 셀 (상태색 배지)


def test_a_row_without_a_deliverable_still_renders():
    items = [row(1, phase_id=1, phase_no=0, milestone_id=10, milestone_no=1, title="제목만")]
    detail = texts_of(open_pptx(build_dashboard_pptx("산출물 없음", items)), 1)
    assert any("(산출물 미지정)" in t for t in detail)


def test_the_legend_lists_owner_kinds_before_statuses():
    """§0.5 범례 개정 — 주관 먼저, 상태 나중."""
    texts = texts_of(open_pptx(build_dashboard_pptx("범례", reference_board())))

    owner_title = next(t for t in texts if t.startswith("주관"))
    status_title = next(t for t in texts if t.startswith("상태"))
    assert texts.index(owner_title) < texts.index(status_title)
    for _key, label, _color in dashboard_theme.OWNER_KINDS:
        assert label in texts
    for status in dashboard_theme.DASH_STATUS_ORDER:
        assert dashboard_theme.STATUS_LABEL[status] in texts


def legend_swatches(deck) -> list[tuple[str, str, str]]:
    """범례 스와치를 `(배경, 테두리, 바)` 로 되읽는다.

    스와치는 겹친 사각형 두 개다 — 같은 y·높이에 왼쪽 끝이 같고 좁은 쪽이 바.
    범례 띠(`LEGEND_Y`) 안에 있는 것만 고른다.
    """
    from app.services.dashboard_pptx_service import LEGEND_Y, SWATCH_BAR_W, SWATCH_W

    EMU = 914400.0
    boxes = [
        (
            round(s.left / EMU, 4), round(s.top / EMU, 4),
            round(s.width / EMU, 4), round(s.height / EMU, 4), s,
        )
        for s in deck.slides[0].shapes
        if s.shape_type is not None and LEGEND_Y <= s.top / EMU <= LEGEND_Y + 0.30
    ]
    bodies = [b for b in boxes if abs(b[2] - SWATCH_W) < 1e-6]
    bars = {(b[0], b[1]): b for b in boxes if abs(b[2] - SWATCH_BAR_W) < 1e-6}

    out = []
    for left, top, _w, _h, body in bodies:
        bar = bars[(left, top)]
        out.append(
            (
                str(body.fill.fore_color.rgb),
                str(body.line.color.rgb),
                str(bar[4].fill.fore_color.rgb),
            )
        )
    return out


def test_the_legend_swatches_are_mini_cards_not_circles():
    """§0.5 범례 개정 — 스와치는 **카드의 축소판**이고, 묶음마다 바뀌는 것이 하나다.

    주관은 좌측 바만, 상태는 배경만 바꾼다. 예전처럼 주관을 원·상태를 사각형으로
    그리면 두 묶음이 서로 다른 어휘처럼 보이고, 색이 카드의 **어디에** 나타나는지
    말해 주지 못한다.
    """
    deck = open_pptx(build_dashboard_pptx("범례", reference_board()))
    swatches = legend_swatches(deck)

    white_bg, white_border, _ = dashboard_theme.DASH_STATUS_STYLE[
        dashboard_theme.DASH_STATUS_ORDER[0]
    ]
    owner_expected = [
        (white_bg, white_border, color) for _key, _label, color in dashboard_theme.OWNER_KINDS
    ]
    status_expected = [
        (
            dashboard_theme.DASH_STATUS_STYLE[s][0],
            dashboard_theme.DASH_STATUS_STYLE[s][1],
            dashboard_theme.UNASSIGNED_OWNER_COLOR,
        )
        for s in dashboard_theme.DASH_STATUS_ORDER
    ]

    assert swatches == owner_expected + status_expected


def test_no_legend_swatch_is_an_oval():
    """원형 스와치로 되돌아가면 여기서 깨진다."""
    from pptx.enum.shapes import MSO_SHAPE

    deck = open_pptx(build_dashboard_pptx("범례", reference_board()))
    from app.services.dashboard_pptx_service import LEGEND_Y

    EMU = 914400.0
    in_legend = [
        s for s in deck.slides[0].shapes
        if s.shape_type is not None and LEGEND_Y <= s.top / EMU <= LEGEND_Y + 0.30
    ]
    def kind(shape):
        # 텍스트 상자는 auto shape 이 아니라 접근 시 ValueError 를 던진다.
        try:
            return shape.auto_shape_type
        except ValueError:
            return None

    assert in_legend, "범례 도형을 찾지 못했다"
    assert all(kind(s) != MSO_SHAPE.OVAL for s in in_legend)
    # 스와치 본체·바는 전부 사각형이다.
    assert sum(1 for s in in_legend if kind(s) == MSO_SHAPE.RECTANGLE) == (
        len(dashboard_theme.OWNER_KINDS) + len(dashboard_theme.DASH_STATUS_ORDER)
    ) * 2


def test_an_empty_board_still_produces_a_valid_status_map():
    """새 프로젝트는 실제로 행이 0개다 — 여기서 터지면 안 된다."""
    deck = open_pptx(build_dashboard_pptx("빈 보드", []))
    assert len(deck.slides) == 1          # 상세 슬라이드는 없다
    texts = texts_of(deck)
    assert any("빈 보드" in t for t in texts)
    assert any(t.startswith("주관") for t in texts)      # 범례는 그대로 나온다


def test_a_board_of_only_gray_rows_renders():
    deck = open_pptx(build_dashboard_pptx("회색뿐", [row(1), row(2)]))
    assert "미배정" in texts_of(deck, 0)
    # 미배정도 자기 상세 슬라이드를 갖는다.
    assert len(deck.slides) == 2
    assert "PROJECT BOARD  ·  UNASSIGNED ITEMS" in texts_of(deck, 1)


# =============================================================================
# 엔드포인트
# =============================================================================
pytestmark_db = pytest.mark.db


@pytest.mark.db
class TestEndpoint:
    @pytest.fixture
    def published(self, db, board):
        from app.models import ItemDocument

        for order in (1, 2):
            item = Item(
                version_id=board.published.id,
                sort_order=order,
                phase_id=board.p0.id,
                milestone_id=board.m01.id,
                title=f"행 {order}",
                deliverable=f"산출물 {order}",
                dash_label=f"라벨 {order}",
            )
            item.owners = [ItemOwner(owner_id=board.o1.id, sort_order=1)]
            item.documents = [ItemDocument(template_document_id=board.d1.id, sort_order=1)]
            db.add(item)
        db.commit()
        return board

    @pytest.fixture
    def project(self, client, published):
        response = client.post(
            f"{API}/projects",
            json={"maker_id": 7, "name": "데모 프로젝트", "template_id": published.wp.id},
        )
        assert response.status_code == 201, response.text
        return response.json()["project"]["id"]

    def test_it_returns_a_pptx(self, client, project):
        response = client.get(f"{API}/projects/{project}/dashboard.pptx")

        assert response.status_code == 200, response.text
        assert response.headers["content-type"] == (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        assert response.content[:2] == b"PK"          # zip 컨테이너

    def test_the_payload_reopens_with_the_board_contents(self, client, project):
        deck = open_pptx(client.get(f"{API}/projects/{project}/dashboard.pptx").content)
        status_map = texts_of(deck, 0)

        # Status Map + Phase 0 상세 한 장 (행 2개).
        assert len(deck.slides) == 2
        assert any("데모 프로젝트" in t for t in status_map)
        assert "1\n라벨 1" in status_map and "2\n라벨 2" in status_map

    def test_the_detail_slide_uses_real_documents_and_owners(self, client, project, db, board):
        """실데이터 경로 — 문서 원문자와 Owner 가 N:M 관계에서 온다."""
        deck = open_pptx(client.get(f"{API}/projects/{project}/dashboard.pptx").content)
        detail = texts_of(deck, 1)

        assert any("산출물 1" in t for t in detail)
        assert "1" in detail                       # board 픽스처의 문서 1번
        assert "DSEP 인프라 담당자" in detail

    def test_the_filename_header_uses_rfc_5987_for_the_korean_name(self, client, project):
        disposition = client.get(
            f"{API}/projects/{project}/dashboard.pptx"
        ).headers["content-disposition"]

        assert disposition.startswith("attachment;")
        # 비ASCII 는 filename* 로. 한글이 헤더에 날것으로 실리면 mojibake 가 된다.
        assert f"filename*=UTF-8''{quote('데모 프로젝트.pptx', safe='')}" in disposition
        # 오래된 클라이언트를 위한 ASCII 폴백도 함께 나간다.
        assert re.search(r'filename="[\x20-\x7e]*"', disposition)

    def test_an_ascii_project_name_is_used_as_is(self, client, published):
        project_id = client.post(
            f"{API}/projects",
            json={"maker_id": 7, "name": "Alpha Board", "template_id": published.wp.id},
        ).json()["project"]["id"]

        disposition = client.get(
            f"{API}/projects/{project_id}/dashboard.pptx"
        ).headers["content-disposition"]
        assert 'filename="Alpha Board.pptx"' in disposition

    def test_a_project_with_no_rows_exports_without_erroring(self, client, board):
        """행 0개 프로젝트 — 500 이 아니라 정상적인 한 장이 나와야 한다."""
        project_id = client.post(
            f"{API}/projects",
            json={"maker_id": 7, "name": "빈 프로젝트", "template_id": board.wp.id},
        ).json()["project"]["id"]

        response = client.get(f"{API}/projects/{project_id}/dashboard.pptx")
        assert response.status_code == 200, response.text
        assert len(open_pptx(response.content).slides) == 1

    def test_a_gray_row_does_not_break_the_export(self, client, project):
        assert client.post(f"{API}/projects/{project}/items").status_code == 201

        response = client.get(f"{API}/projects/{project}/dashboard.pptx")
        assert response.status_code == 200, response.text
        deck = open_pptx(response.content)
        assert "미배정" in texts_of(deck, 0)
        # Status Map + Phase 0 상세 + 미배정 상세.
        assert len(deck.slides) == 3

    def test_an_unknown_project_is_404(self, client):
        assert client.get(f"{API}/projects/999999/dashboard.pptx").status_code == 404

    def test_the_export_is_a_read_only_route(self, client):
        schema = client.app.openapi()
        assert set(schema["paths"]["/api/v1/projects/{project_id}/dashboard.pptx"]) == {"get"}


# =============================================================================
# 색 상수가 프론트와 어긋나지 않는가 (§0.5.6 — "어긋나면 스펙 위반")
# =============================================================================
def frontend_theme_source() -> str:
    if not FRONTEND_THEME.exists():
        pytest.skip(f"프론트 테마 파일이 없습니다: {FRONTEND_THEME}")
    return FRONTEND_THEME.read_text(encoding="utf-8")


def test_phase_palette_matches_the_frontend():
    source = frontend_theme_source()
    declared = re.search(r"DASH_PHASE_COLORS\s*=\s*\[([^\]]+)\]", source)
    assert declared, "프론트에서 DASH_PHASE_COLORS 를 찾지 못했다"

    frontend = [h.upper() for h in re.findall(r"#([0-9a-fA-F]{6})", declared.group(1))]
    assert frontend == list(dashboard_theme.DASH_PHASE_COLORS)


PLAN = REPO_ROOT / "plan.md"


def test_status_colours_match_the_plan_md_table():
    """**정본은 plan.md §0.5 의 표다.** 그 표를 직접 파싱해 대조한다.

    프론트 대조(`test_status_colours_match_the_frontend`)만으로는 부족하다 —
    두 파일이 **같이** 틀리면 서로 일치하므로 통과한다. 스펙에 못을 박아야
    "둘이 합의했지만 사용자가 정한 값이 아닌" 상태를 잡을 수 있다.
    """
    source = PLAN.read_text(encoding="utf-8")
    table: dict[str, tuple[str, str, str]] = {}
    for line in source.splitlines():
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if len(cells) != 4:
            continue
        hexes = [re.fullmatch(r"`#([0-9a-fA-F]{6})`", c) for c in cells[1:]]
        if not all(hexes):
            continue
        name = re.search(r"\b(NOT_STARTED|IN_PROGRESS|DONE|HOLD|NA)\b", cells[0])
        if name:
            table[name.group(1)] = tuple(h.group(1).upper() for h in hexes)

    assert len(table) == 5, f"plan.md §0.5 상태 표를 읽지 못했다: {sorted(table)}"

    ours = {status.value: value for status, value in dashboard_theme.DASH_STATUS_STYLE.items()}
    assert table == ours

    # 2026-08-08 개정분을 문자 그대로 못박는다 — 표 파싱이 느슨해져도 이건 남는다.
    assert ours["IN_PROGRESS"] == ("D1FAE5", "34D399", "065F46")
    assert ours["DONE"] == ("CBD5E1", "94A3B8", "1E293B")
    assert ours["NA"] == ("334155", "1E293B", "CBD5E1")
    # 진행전·보류는 유지.
    assert ours["NOT_STARTED"] == ("FFFFFF", "94A3B8", "334155")
    assert ours["HOLD"] == ("FEE2E2", "FCA5A5", "991B1B")


def _relative_luminance(hex_color: str) -> float:
    def channel(value: int) -> float:
        c = value / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

    r, g, b = (int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
    return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)


def contrast_ratio(fg: str, bg: str) -> float:
    a, b = _relative_luminance(fg), _relative_luminance(bg)
    lighter, darker = max(a, b), min(a, b)
    return (lighter + 0.05) / (darker + 0.05)


def test_every_status_keeps_its_text_readable_on_its_own_background():
    """NA 가 **짙은 배경 + 밝은 글자**로 뒤집혔다 — 대비를 눈이 아니라 수치로 지킨다.

    카드 문구도 상세 슬라이드의 No 배지도 배경이 상태색이고 글자가 상태 글자색이라,
    한 쌍이라도 어긋나면 그 상태의 카드가 통째로 읽히지 않는다. WCAG AA 본문
    기준 4.5:1 을 하한으로 둔다.
    """
    for status, (bg, _border, text) in dashboard_theme.DASH_STATUS_STYLE.items():
        ratio = contrast_ratio(text, bg)
        assert ratio >= 4.5, f"{status.value}: 글자 대비가 {ratio:.2f}:1 로 낮다"


def test_the_na_card_is_light_text_on_a_dark_ground():
    """방향 자체를 못박는다 — 흐림 처리로 되돌아가면 여기서 깨진다."""
    bg, _border, text = dashboard_theme.DASH_STATUS_STYLE[ItemStatus.NA]
    assert _relative_luminance(bg) < _relative_luminance(text)


def test_status_colours_match_the_frontend():
    source = frontend_theme_source()
    block = re.search(r"DASH_STATUS_STYLE[^{]*\{(.+?)\n\}", source, re.S)
    assert block, "프론트에서 DASH_STATUS_STYLE 을 찾지 못했다"

    frontend: dict[str, tuple[str, str, str]] = {}
    for name, body in re.findall(r"(\w+):\s*\{([^}]*)\}", block.group(1)):
        fields = dict(re.findall(r"(\w+):\s*'#([0-9a-fA-F]{6})'", body))
        frontend[name] = (
            fields["bg"].upper(), fields["border"].upper(), fields["text"].upper()
        )

    ours = {status.value: value for status, value in dashboard_theme.DASH_STATUS_STYLE.items()}
    assert frontend == ours


def test_owner_kind_colours_and_labels_match_the_frontend():
    source = frontend_theme_source()
    block = re.search(r"OWNER_KINDS\s*=\s*\[(.+?)\]\s*as const", source, re.S)
    assert block, "프론트에서 OWNER_KINDS 를 찾지 못했다"

    frontend = [
        (key, label, color.upper())
        for key, label, color in re.findall(
            r"key:\s*'(\w+)',\s*label:\s*'([^']*)',\s*color:\s*'#([0-9a-fA-F]{6})'",
            block.group(1),
        )
    ]
    assert frontend == [
        (key, label, color) for key, label, color in dashboard_theme.OWNER_KINDS
    ]


def test_the_unassigned_band_colour_matches_the_frontend():
    source = frontend_theme_source()
    declared = re.search(r"DASH_UNASSIGNED\s*=\s*'#([0-9a-fA-F]{6})'", source)
    assert declared
    assert declared.group(1).upper() == dashboard_theme.DASH_UNASSIGNED


# =============================================================================
# ownerKind 휴리스틱 — 프론트와 같은 규칙
# =============================================================================
@pytest.mark.parametrize(
    "names,expected",
    [
        ([], "NONE"),
        (None, "NONE"),
        (["사내 개발부서"], "INTERNAL_DEV"),
        (["DSEP 인프라 담당자"], "DSEP"),
        (["설비사"], "MAKER"),
        (["설비사 PM"], "MAKER"),
        (["공동(구매·법무·보안)"], "JOINT"),
        (["가+나"], "JOINT"),
        # 2명 이상은 이름과 무관하게 공동 — 추측이 아니라 구조적 사실이다.
        (["DSEP 인프라 담당자", "설비사"], "JOINT"),
        (["사내 IT·보안"], "INTERNAL_DEV"),
        # ⚠️ **판정 순서를 가르는 유일한 케이스.** 두 단어를 다 가진 이름은
        # 먼저 검사하는 쪽이 이긴다. 이 줄이 없으면 위 케이스들은 DSEP 검사와
        # 설비사 검사를 맞바꿔도 전부 통과한다 — 실제로 그렇게 통과시켜 보고
        # 추가했다 (HANDOFF §5.1d: "이 픽스처가 애초에 실패할 수 있는가").
        (["DSEP 설비사 지원 담당"], "DSEP"),
        # 공동이 가장 먼저인 것도 같은 방식으로 고정한다.
        (["공동 DSEP 설비사"], "JOINT"),
    ],
)
def test_owner_kind_matches_the_frontend_heuristic(names, expected):
    assert dashboard_theme.owner_kind(names) == expected


def test_the_heuristic_checks_in_the_same_order_as_the_frontend():
    """순서가 규칙의 일부다 — 두 단어를 가진 이름의 답이 순서로 갈린다.

    프론트 소스에서 판정 순서를 뽑아 백엔드 구현이 **실제로 그 순서로 답하는지**
    확인한다. 소스 문자열만 비교하면 프론트 안에서만 닫힌 검사가 되어, 백엔드가
    갈라져도 통과한다.
    """
    source = frontend_theme_source()
    order = re.findall(r"if \(name\.includes\('([^']+)'\)", source)
    assert order[:3] == ["공동", "DSEP", "설비사"]

    # 앞선 검사가 뒤선 검사를 이긴다 — 두 단어를 모두 가진 이름으로 확인.
    assert dashboard_theme.owner_kind(["DSEP 설비사"]) == "DSEP"
    assert dashboard_theme.owner_kind(["공동 DSEP"]) == "JOINT"
    assert dashboard_theme.owner_kind(["공동 설비사"]) == "JOINT"
