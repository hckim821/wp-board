"""대시보드 → 다중 슬라이드 PPTX — plan.md §0.5.6 (다중 슬라이드로 격상).

내보내기는 백엔드가 만든다 (CLAUDE.md 의 openpyxl 경로와 동형). 화면 캡처가 아니라
데이터에서 다시 그리는 방식이라, 스크롤 밖의 항목도 빠지지 않고 폰트·색이 실행
환경에 좌우되지 않는다.

## 시각 정본은 `docs/DSEP_AI_Project_Board_Guide.pptx` 다

아래 치수·폰트 크기·크롬 색은 **그 덱을 python-pptx 로 파싱해 뽑은 실측값**이다
(추정이 아니다). 슬라이드 1 은 그 덱의 slide 1(Status Map), 상세 슬라이드는
slide 2~8 의 형식을 따른다.

한 가지는 정본 덱을 따르지 않는다: **Phase·상태·주관 팔레트**. 그건 plan.md §0.5
가 못박은 값이고 화면(`frontend/src/theme/dashboard.ts`)과 픽셀 단위로 맞아야 하며,
드리프트를 테스트가 잡고 있다. 정본 덱의 색은 §0.5 와 사실상 같은 색이지만
(`8E7CC3` vs `8F7CC3`) 한 자리씩 다르므로, 그쪽을 쓰면 화면과 어긋난다. 제목·
테두리·푸터처럼 **화면에 대응물이 없는 색만** 덱에서 가져왔다.

## 구성

* **슬라이드 1 — Status Map**: Phase chevron + 마일스톤 헤더 + 카드 격자 + 범례.
  언제나 한 장이며 컬럼 수·최대 카드 수로 자동 축소한다.
* **슬라이드 2+ — Phase 별 상세**: Phase 하나가 한 장 이상. 한 장에 5행씩 끊고,
  나뉘면 제목에 `(1/2)` 가 붙는다. 미배정 행은 맨 뒤에 자기 그룹으로 나온다.

## python-pptx 는 늦게 import 한다

모듈 최상단에서 import 하면 `create_wp_router()` 가 python-pptx 를 하드 의존하게
된다. 함수 안에서 import 하고, 없으면 이 엔드포인트만 501 로 답한다.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO

from sqlalchemy.orm import Session

from ..core.exceptions import WpError
from ..schemas.item import ItemOut
from .dashboard_theme import (
    DASH_STATUS_ORDER,
    DASH_STATUS_STYLE,
    DASH_UNASSIGNED,
    DECK_BODY,
    DECK_BORDER,
    DECK_EYEBROW,
    DECK_INK,
    DECK_MUTED,
    OWNER_KINDS,
    STATUS_LABEL,
    UNASSIGNED_OWNER_COLOR,
    dashboard_text,
    owner_color,
    phase_color,
    tint,
)


class ExportUnavailableError(WpError):
    """python-pptx 미설치. **501** — 요청이 틀린 게 아니라 서버가 못 하는 것이다."""

    status_code = 501
    code = "EXPORT_UNAVAILABLE"


# =============================================================================
# 치수 — 정본 덱 실측값 (inch)
# =============================================================================
SLIDE_W = 13.333
SLIDE_H = 7.5

MARGIN_X = 0.52
#: 정본 덱의 여백은 **좌우 비대칭**이다 — 13 컬럼이 `0.52 + 13×0.95 = 12.87` 까지
#: 차고 오른쪽에 0.46 이 남는다. 좌우를 같다고 두면 기준 케이스조차 피치가
#: 0.9456 으로 줄어 덱과 미묘하게 어긋난다 (테스트가 그걸 잡았다).
MARGIN_RIGHT = 0.46
EYEBROW_Y, EYEBROW_H, EYEBROW_PT = 0.38, 0.26, 10.5
TITLE_Y, TITLE_H, TITLE_PT = 0.62, 0.62, 24.0
FOOTER_Y, FOOTER_H, FOOTER_PT = 7.08, 0.30, 8.5

# --- 슬라이드 1 (Status Map) --------------------------------------------------
CHEVRON_Y, CHEVRON_H, PHASE_PT = 1.38, 0.32, 9.5
#: chevron 은 컬럼보다 살짝 왼쪽에서 시작하고, 그 안에 라벨이 들어간다.
CHEVRON_INSET = 0.15
MSHEAD_Y, MSHEAD_H, MSHEAD_PT = 1.76, 0.52, 8.0
CARDS_Y, CARD_H, CARD_GAP, CARD_PT = 2.36, 0.82, 0.07, 8.5
COL_PITCH, COL_W, OWNER_BAR_W = 0.95, 0.90, 0.05
LEGEND_Y, LEGEND_PT = 6.18, 8.5
#: 범례 스와치 = **카드의 축소판** (§0.5 범례 개정, 화면의 `LegendSwatch.vue`).
#: 웹의 22×14px 비율을 그대로 옮겼다. 예전처럼 주관을 원, 상태를 정사각형으로
#: 그리면 두 묶음이 서로 다른 어휘처럼 보이고, 무엇보다 **색이 카드의 어디에
#: 나타나는지**를 말해 주지 못한다 — 그러면 "좌측 바 = 주관" 을 글로 덧붙여야 한다.
SWATCH_W, SWATCH_H, SWATCH_BAR_W = 0.26, 0.165, 0.05
#: 카드가 침범하면 안 되는 하한 (범례 위).
CARDS_BOTTOM = 6.05

# --- 슬라이드 2+ (Phase 별 상세) ----------------------------------------------
COLHEAD_Y, COLHEAD_H, COLHEAD_PT = 1.55, 0.26, 9.0
ROW_Y0, ROW_H, ROW_PITCH = 1.95, 0.82, 0.90
ROW_X, ROW_W = 0.55, 12.23
PHASE_BAR_W = 0.07
NO_X, NO_W, NO_PT = 0.68, 0.50, 13.0
MAIN_X, MAIN_W, MAIN_H, MAIN_PT = 1.30, 7.10, 0.70, 10.5
DOC_X, DOC_W, DOC_PT = 8.50, 1.00, 12.0
OWNER_X, OWNER_W, OWNER_PT = 9.60, 3.15, 8.5

#: 한 장에 담는 행 수. 정본 덱의 상세 슬라이드가 4~5행이고, 5행이면
#: `1.95 + 5×0.90 = 6.45` 로 푸터(7.08) 위에 넉넉히 들어간다.
ROWS_PER_SLIDE = 5


# =============================================================================
# 슬라이드 1 레이아웃 — 순수 계산 (python-pptx 없이 테스트된다)
# =============================================================================
@dataclass
class Card:
    no: int
    text: str
    status: object
    owners: list[str]


@dataclass
class Column:
    """마일스톤 하나 = 컬럼 하나. 미배정도 컬럼 하나로 취급한다."""

    number_label: str
    name: str
    cards: list[Card] = field(default_factory=list)


@dataclass
class Band:
    """Phase 밴드. 자기 마일스톤 컬럼들을 가로로 덮는다."""

    label: str
    color: str
    columns: list[Column] = field(default_factory=list)


def build_layout(items: list[ItemOut]) -> list[Band]:
    """행 목록 → Phase 밴드 → 마일스톤 컬럼 → 카드.

    `composables/useDashboard.ts:buildDashboardLayout` 과 **같은 규칙**이다.

    * 회색 행(미배정)은 격자에서 빠져 **맨 뒤 무색 밴드**로 모인다. 연속성
      판정에서는 투명하지만(§0.2.1), 대시보드의 컬럼은 소속의 주장이므로 소속이
      없는 행은 어느 컬럼에도 들어가지 않는다.
    * 조각난 Phase 도 **밴드 하나**다. 첫 등장이 순서를 정하고 뒤에 다시 나오는
      같은 Phase 는 그 밴드에 합쳐진다 — 같은 이름표를 단 밴드 둘은 조각남을
      알리는 더 나쁜 방법이고, 그건 발행 검증(V4/V5)이 할 일이다.
    * Phase 는 있는데 Milestone 이 없는 행은 그 Phase 안의 **이름 없는 버킷**에
      담긴다. 두 계층 모두에서 합법이므로 떨어뜨릴 수 없다.
    """
    bands: list[Band] = []
    by_phase: dict[int, Band] = {}
    by_phase_columns: dict[int, dict[str, Column]] = {}
    unassigned = Column(number_label="", name="미배정")

    for item in items:
        card = Card(
            no=item.row_no,
            text=dashboard_text(item.dash_label, item.deliverable, item.title),
            status=item.status,
            owners=[o.name for o in item.owners],
        )

        if item.phase_id is None:
            unassigned.cards.append(card)
            continue

        band = by_phase.get(item.phase_id)
        if band is None:
            label = (
                f"Phase {item.phase_no}. {item.phase_name or ''}".strip()
                if item.phase_no is not None
                else (item.phase_name or "")
            )
            band = Band(label=label, color=phase_color(item.phase_no))
            by_phase[item.phase_id] = band
            by_phase_columns[item.phase_id] = {}
            bands.append(band)

        key = "none" if item.milestone_id is None else str(item.milestone_id)
        column = by_phase_columns[item.phase_id].get(key)
        if column is None:
            number = (
                f"{item.phase_no}.{item.milestone_no}"
                if item.phase_no is not None and item.milestone_no is not None
                else ""
            )
            column = Column(
                number_label=number,
                name="미지정" if item.milestone_id is None else (item.milestone_name or ""),
            )
            by_phase_columns[item.phase_id][key] = column
            band.columns.append(column)
        column.cards.append(card)

    if unassigned.cards:
        bands.append(Band(label="미배정", color=DASH_UNASSIGNED, columns=[unassigned]))

    return bands


@dataclass(frozen=True)
class Metrics:
    """슬라이드 1 의 자동 축소 결과. **카드 높이는 하나뿐이다** (§0.5.4b)."""

    column_pitch: float
    column_w: float
    card_h: float
    card_pitch: float
    phase_pt: float
    milestone_pt: float
    card_pt: float


def compute_metrics(bands: list[Band]) -> Metrics:
    """컬럼 수(가로)와 최대 카드 수(세로)만 보고 한 장에 맞춘다.

    정본 덱의 값(피치 0.95 · 카드 0.82)이 **상한**이다. 기준 케이스(마일스톤 13개
    · 최대 4장)에서는 그대로 쓰이고, 더 빽빽해질 때만 줄어든다. 폰트에는 하한을
    둔다 — 하한에 걸리면 글자가 잘릴 수 있는데, 잘리는 것이 다음 장으로 밀려나는
    것보다 낫다 (화면에서도 hover 팝오버가 보완하는 부분이다, §0.5.4b).
    """
    total_columns = max(sum(len(b.columns) for b in bands), 1)
    tallest = max((len(c.cards) for b in bands for c in b.columns), default=0)

    usable_w = SLIDE_W - MARGIN_X - MARGIN_RIGHT
    column_pitch = min(COL_PITCH, usable_w / total_columns)
    column_w = column_pitch - (COL_PITCH - COL_W)

    # 세로는 **피치를 먼저** 정하고 거기서 간격을 떼어 낸다. 카드 높이에 하한을
    # 두고 간격을 고정하면, 카드가 많을 때 하한이 이겨서 범례를 덮어 버린다
    # (실제로 20장짜리 컬럼에서 그렇게 넘쳤다). 간격이 카드와 함께 줄면 그 일이
    # 생기지 않는다.
    available_h = CARDS_BOTTOM - CARDS_Y
    if tallest <= 0:
        card_pitch, card_h = CARD_H + CARD_GAP, CARD_H
    else:
        card_pitch = min(CARD_H + CARD_GAP, available_h / tallest)
        card_h = max(card_pitch - min(CARD_GAP, card_pitch * 0.08), 0.10)

    scale = min(1.0, column_pitch / COL_PITCH)
    return Metrics(
        column_pitch=column_pitch,
        column_w=column_w,
        card_h=card_h,
        card_pitch=card_pitch,
        phase_pt=max(6.5, PHASE_PT * scale),
        milestone_pt=max(5.5, MSHEAD_PT * scale),
        card_pt=max(5.5, CARD_PT * scale),
    )


# =============================================================================
# 슬라이드 2+ 레이아웃 — Phase 별 상세, 5행씩
# =============================================================================
@dataclass
class DetailRow:
    no: int
    deliverable: str
    milestone: str
    title: str
    documents: list[str]
    owners: list[str]
    status: object


@dataclass
class DetailPage:
    """상세 슬라이드 한 장."""

    #: `Phase 0. Pre-Infrastructure Setup` / `미배정`
    group_label: str
    #: 눈썹줄 오른쪽 문구 (`PHASE 0 ITEMS`).
    eyebrow: str
    #: Phase 색 (행 왼쪽 세로 바).
    color: str
    page_index: int
    page_count: int
    rows: list[DetailRow]

    @property
    def title(self) -> str:
        """`Phase 0. … (1/2) — 항목 1~5`.

        나뉘지 않았으면 `(1/1)` 을 붙이지 않는다 — 한 장뿐인데 페이지 표기가
        붙으면 뒤에 더 있는 것처럼 읽힌다.
        """
        span = f"항목 {self.rows[0].no}~{self.rows[-1].no}" if self.rows else "항목 없음"
        if self.page_count > 1:
            return f"{self.group_label} ({self.page_index}/{self.page_count}) — {span}"
        return f"{self.group_label} — {span}"


def build_detail_pages(items: list[ItemOut]) -> list[DetailPage]:
    """Phase 별로 묶고 한 장 `ROWS_PER_SLIDE` 행씩 끊는다.

    묶는 규칙은 슬라이드 1 과 같다 — 첫 등장 순서, 조각난 Phase 는 한 묶음,
    **미배정은 맨 뒤 자기 그룹**. 두 슬라이드가 같은 순서로 읽혀야 하므로
    여기서만 다르게 묶을 이유가 없다.
    """
    groups: list[tuple[str, str, str, list[DetailRow]]] = []
    by_phase: dict[int, int] = {}           # phase_id → groups 인덱스
    unassigned: list[DetailRow] = []

    for item in items:
        row = DetailRow(
            no=item.row_no,
            deliverable=(item.deliverable or "").strip(),
            milestone=" ".join(
                part for part in (item.milestone_no_display, item.milestone_name) if part
            ).strip(),
            title=(item.title or "").strip(),
            # §0.5.10 — 원문자가 아니라 **파생 표시 번호**다. 꺼진 문서는
            # `build_item_views` 단계에서 이미 빠져 있으므로 여기 `no` 는 항상 정수다.
            documents=[str(d.no) for d in item.documents],
            owners=[o.name for o in item.owners],
            status=item.status,
        )

        if item.phase_id is None:
            unassigned.append(row)
            continue

        index = by_phase.get(item.phase_id)
        if index is None:
            label = (
                f"Phase {item.phase_no}. {item.phase_name or ''}".strip()
                if item.phase_no is not None
                else (item.phase_name or "Phase")
            )
            eyebrow = (
                f"PHASE {item.phase_no} ITEMS" if item.phase_no is not None else "PHASE ITEMS"
            )
            by_phase[item.phase_id] = len(groups)
            groups.append((label, eyebrow, phase_color(item.phase_no), []))
            index = by_phase[item.phase_id]
        groups[index][3].append(row)

    if unassigned:
        groups.append(("미배정", "UNASSIGNED ITEMS", DASH_UNASSIGNED, unassigned))

    pages: list[DetailPage] = []
    for label, eyebrow, color, rows in groups:
        chunks = [
            rows[i : i + ROWS_PER_SLIDE] for i in range(0, len(rows), ROWS_PER_SLIDE)
        ] or [[]]
        for page_index, chunk in enumerate(chunks, start=1):
            pages.append(
                DetailPage(
                    group_label=label,
                    eyebrow=eyebrow,
                    color=color,
                    page_index=page_index,
                    page_count=len(chunks),
                    rows=chunk,
                )
            )
    return pages


# =============================================================================
# 생성
# =============================================================================
def build_dashboard_pptx(project_name: str, items: list[ItemOut]) -> bytes:
    """Status Map 한 장 + Phase 별 상세 여러 장.

    **행이 0개여도 정상적으로 만들어진다** — Status Map 한 장만 나온다. 빈 보드는
    실제로 존재하는 상태(새 프로젝트)이므로 여기서 터지면 안 된다.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - 설치된 환경에서는 도달하지 않는다
        raise ExportUnavailableError(
            "PPT 내보내기를 사용하려면 서버에 python-pptx 가 설치되어 있어야 합니다.",
            detail={"package": "python-pptx"},
        ) from exc

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    blank = presentation.slide_layouts[6]      # 자리표시자 없는 레이아웃

    def new_slide():
        return presentation.slides.add_slide(blank)

    def textbox(slide, left, top, width, height, runs, *, size, bold=True,
                color=DECK_BODY, align=PP_ALIGN.CENTER):
        """`runs` 는 문자열이거나 `(텍스트, 굵게, 색)` 튜플의 목록."""
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.03)
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align

        pieces = [(runs, bold, color)] if isinstance(runs, str) else runs
        for text, is_bold, fg in pieces:
            run = paragraph.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = is_bold
            run.font.color.rgb = RGBColor.from_string(fg)
        return box

    def card_textbox(slide, left, top, width, height, no, text, *, size, color):
        """카드 문구 — **두 문단**이다 (§0.5-2 개정).

        첫 줄은 No 만 가운데·굵게, 둘째 줄은 라벨을 왼쪽으로. 웹 대시보드 카드와
        같은 포맷이며, 한 문단에 `1 | 라벨` 로 합치면 번호가 문장의 일부처럼
        읽혀 눈이 번호를 먼저 잡지 못한다.

        정렬이 문단마다 다르므로 일반 `textbox` 를 쓸 수 없다 — 정렬은 run 이
        아니라 문단의 속성이다.
        """
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.03)
        frame.margin_top = frame.margin_bottom = 0
        # **위 정렬.** 가운데 앵커로 두면 라벨이 한 줄인 카드와 두 줄로 접히는
        # 카드에서 번호의 높이가 서로 달라져, 격자를 훑는 눈이 번호 줄을 잃는다.
        # 카드 높이는 이미 전 보드에서 하나로 통일돼 있으므로(§0.5.4b), 위를
        # 기준으로 맞추면 모든 번호가 같은 선에서 시작한다.
        frame.vertical_anchor = MSO_ANCHOR.TOP

        head = frame.paragraphs[0]
        head.alignment = PP_ALIGN.CENTER
        run = head.add_run()
        run.text = str(no)
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(color)

        if text:
            body = frame.add_paragraph()
            body.alignment = PP_ALIGN.LEFT
            run = body.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = False
            run.font.color.rgb = RGBColor.from_string(color)
        return box

    def shape(slide, kind, left, top, width, height, fill, *, line=None):
        item = slide.shapes.add_shape(
            kind, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        item.fill.solid()
        item.fill.fore_color.rgb = RGBColor.from_string(fill)
        if line is None:
            item.line.fill.background()
        else:
            item.line.color.rgb = RGBColor.from_string(line)
            item.line.width = Pt(0.75)
        item.shadow.inherit = False
        return item

    def box(slide, left, top, width, height, fill, *, line=None):
        return shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill, line=line)

    def chrome(slide, eyebrow: str, title: str):
        textbox(slide, MARGIN_X, EYEBROW_Y, SLIDE_W - MARGIN_X * 2, EYEBROW_H,
                eyebrow, size=EYEBROW_PT, color=DECK_EYEBROW)
        textbox(slide, MARGIN_X, TITLE_Y, SLIDE_W - MARGIN_X * 2, TITLE_H,
                title, size=TITLE_PT, color=DECK_INK)

    def footer(slide, page_no: int):
        textbox(slide, ROW_X, FOOTER_Y, 8.0, FOOTER_H, project_name,
                size=FOOTER_PT, bold=False, color=DECK_MUTED, align=PP_ALIGN.LEFT)
        textbox(slide, SLIDE_W - MARGIN_X - 0.88, FOOTER_Y, 0.88, FOOTER_H, str(page_no),
                size=FOOTER_PT, bold=False, color=DECK_MUTED, align=PP_ALIGN.RIGHT)

    # =========================================================================
    # 슬라이드 1 — Status Map
    # =========================================================================
    bands = build_layout(items)
    m = compute_metrics(bands)
    slide = new_slide()
    chrome(slide, "PROJECT BOARD  ·  STATUS MAP", f"{project_name}  ({len(items)}항목)")

    x = MARGIN_X
    for band in bands:
        span = max(len(band.columns), 1) * m.column_pitch
        # Phase chevron — 왼쪽으로 살짝 내밀어 진행 방향을 보인다.
        shape(slide, MSO_SHAPE.CHEVRON, x, CHEVRON_Y, span - (COL_PITCH - COL_W),
              CHEVRON_H, band.color)
        textbox(slide, x + CHEVRON_INSET, CHEVRON_Y, span - CHEVRON_INSET * 2, CHEVRON_H,
                band.label, size=m.phase_pt, color="FFFFFF")

        head_fill = tint(band.color)
        cx = x
        for column in band.columns:
            box(slide, cx, MSHEAD_Y, m.column_w, MSHEAD_H, head_fill)
            label = "\n".join(part for part in (column.number_label, column.name) if part)
            textbox(slide, cx + 0.04, MSHEAD_Y, m.column_w - 0.08, MSHEAD_H,
                    label, size=m.milestone_pt, color=DECK_INK)

            cy = CARDS_Y
            for card in column.cards:
                fill, border, text_color = DASH_STATUS_STYLE[card.status]
                box(slide, cx, cy, m.column_w, m.card_h, fill, line=border)
                # 주관 = 좌측 세로 바.
                box(slide, cx, cy, OWNER_BAR_W, m.card_h, owner_color(card.owners))
                card_textbox(slide, cx + 0.11, cy + 0.03, m.column_w - 0.16,
                             m.card_h - 0.06, card.no, card.text,
                             size=m.card_pt, color=text_color)
                cy += m.card_pitch
            cx += m.column_pitch
        x += span

    # --- 범례 (주관 먼저, 상태 나중 — §0.5 범례 개정) -------------------------
    #
    # 정본 덱은 상태를 먼저 놓지만 §0.5 개정이 순서를 뒤집었고, 화면 범례도 그렇다.
    # 지오메트리는 덱에서, 순서는 스펙에서 가져온다.
    def swatch(left: float, bg: str, border: str, bar: str):
        """카드의 축소판. **바뀌는 것 하나만** 바뀐다 — 주관은 바, 상태는 배경."""
        top = LEGEND_Y + (0.30 - SWATCH_H) / 2
        box(slide, left, top, SWATCH_W, SWATCH_H, bg, line=border)
        box(slide, left, top, SWATCH_BAR_W, SWATCH_H, bar)

    def legend_entry(left: float, label: str, bg: str, border: str, bar: str) -> float:
        swatch(left, bg, border, bar)
        textbox(slide, left + SWATCH_W + 0.06, LEGEND_Y, 1.5, 0.30, label,
                size=LEGEND_PT, bold=False, color=DECK_MUTED, align=PP_ALIGN.LEFT)
        return left + SWATCH_W + 0.06 + len(label) * 0.115 + 0.20

    # 주관 스와치는 **배경을 진행전(흰색)으로 고정**하고 좌측 바만 바꾼다.
    neutral_bg, neutral_border, _ = DASH_STATUS_STYLE[DASH_STATUS_ORDER[0]]
    lx = MARGIN_X
    textbox(slide, lx, LEGEND_Y, 1.2, 0.30, "주관 (좌측 바)",
            size=LEGEND_PT, color=DECK_INK, align=PP_ALIGN.LEFT)
    lx += 1.25
    for _key, label, color in OWNER_KINDS:
        lx = legend_entry(lx, label, neutral_bg, neutral_border, color)

    # 상태 스와치는 반대로 **좌측 바를 미지정 slate 로 고정**하고 배경만 바꾼다.
    lx = max(lx, 7.5)
    textbox(slide, lx, LEGEND_Y, 0.9, 0.30, "상태 (배경)",
            size=LEGEND_PT, color=DECK_INK, align=PP_ALIGN.LEFT)
    lx += 0.95
    for status in DASH_STATUS_ORDER:
        fill, border, _text = DASH_STATUS_STYLE[status]
        lx = legend_entry(lx, STATUS_LABEL[status], fill, border, UNASSIGNED_OWNER_COLOR)

    footer(slide, 1)

    # =========================================================================
    # 슬라이드 2+ — Phase 별 상세
    # =========================================================================
    for page_no, page in enumerate(build_detail_pages(items), start=2):
        slide = new_slide()
        chrome(slide, f"PROJECT BOARD  ·  {page.eyebrow}", page.title)

        textbox(slide, MAIN_X - 0.10, COLHEAD_Y, MAIN_W + 0.10, COLHEAD_H,
                "산출물 (Deliverable)  ·  Milestone  /  Key Action",
                size=COLHEAD_PT, color=DECK_MUTED)
        textbox(slide, DOC_X, COLHEAD_Y, DOC_W, COLHEAD_H, "문서",
                size=COLHEAD_PT, color=DECK_MUTED)
        textbox(slide, OWNER_X, COLHEAD_Y, OWNER_W, COLHEAD_H, "Owner",
                size=COLHEAD_PT, color=DECK_MUTED)

        y = ROW_Y0
        for row in page.rows:
            fill, border, text_color = DASH_STATUS_STYLE[row.status]
            box(slide, ROW_X, y, ROW_W, ROW_H, "FFFFFF", line=DECK_BORDER)
            # 왼쪽 세로 바 = Phase 색 (정본 덱과 같다).
            box(slide, ROW_X, y, PHASE_BAR_W, ROW_H, page.color)
            # **상태 색 확장** (§0.5.6 개정) — 정본 덱에는 없다. No 셀을 상태색
            # 배지로 칠해, 상세 슬라이드에서도 진행 상태가 한눈에 들어오게 한다.
            box(slide, NO_X, y + 0.08, NO_W, ROW_H - 0.16, fill, line=border)
            textbox(slide, NO_X, y + 0.08, NO_W, ROW_H - 0.16, str(row.no),
                    size=NO_PT, color=text_color)

            tail = " ".join(part for part in (
                f"·  {row.milestone}" if row.milestone else "",
                f"| {row.title}" if row.title else "",
            ) if part)
            runs = [(row.deliverable or "(산출물 미지정)", True, DECK_BODY)]
            if tail:
                runs.append((f"   {tail}", False, DECK_MUTED))
            textbox(slide, MAIN_X, y + 0.06, MAIN_W, MAIN_H, runs,
                    size=MAIN_PT, color=DECK_BODY)

            textbox(slide, DOC_X, y, DOC_W, ROW_H, " ".join(row.documents),
                    size=DOC_PT, color=DECK_INK)
            textbox(slide, OWNER_X, y, OWNER_W, ROW_H, "+".join(row.owners),
                    size=OWNER_PT, bold=False, color=DECK_BODY)
            y += ROW_PITCH

        footer(slide, page_no)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def dashboard_pptx(session: Session, project_id: int) -> tuple[str, bytes]:
    """`(프로젝트명, 파일 바이트)`. 그리드와 **같은 조회 경로**를 쓴다.

    `build_item_views` 를 거치므로 번호·소속·문서·Owner 이름이 화면과 동일하다 —
    내보내기가 자기만의 조회를 갖고 있으면 두 결과가 조용히 갈린다.

    """
    from . import item_service, project_service

    project = project_service.get_project(session, project_id)
    board = project_service.board_of(project)
    items = item_service.load_ordered_items(session, board)
    views = item_service.build_item_views(session, board, items)
    return project.name, build_dashboard_pptx(project.name, views)
